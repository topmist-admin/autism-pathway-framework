#!/usr/bin/env python3
"""
Generate a detailed PowerPoint presentation for the Autism Pathway Framework.
Covers tech stack and ML algorithms (implemented and planned).
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# Create presentation with widescreen dimensions
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Color scheme - Autism Awareness Theme
# Based on official autism awareness colors from https://www.schemecolor.com/autism.php
# Combined with modern scientific/research aesthetic

# Primary: Blue Cola (Autism Awareness Blue) - represents serenity, knowledge
PRIMARY_COLOR = RGBColor(0x0C, 0x96, 0xE4)  # #0C96E4 - Blue Cola

# Secondary: American Green (Autism Awareness) - represents growth, hope
SECONDARY_COLOR = RGBColor(0x1E, 0xB7, 0x42)  # #1EB742 - American Green

# Accent: Deep Carmine Pink (Autism Awareness) - represents energy, awareness
ACCENT_COLOR = RGBColor(0xFF, 0x33, 0x34)  # #FF3334 - Deep Carmine Pink

# Warm Yellow (Autism Awareness) - for highlights
HIGHLIGHT_COLOR = RGBColor(0xF6, 0xBB, 0x00)  # #F6BB00 - American Yellow

# Text and background colors
DARK_TEXT = RGBColor(0x1f, 0x29, 0x37)  # Dark slate for readability
LIGHT_TEXT = RGBColor(0x4b, 0x55, 0x63)  # Medium gray for subtitles
LIGHT_BG = RGBColor(0xf0, 0xf9, 0xff)  # Light blue tint background

# Status colors using autism palette
SUCCESS_COLOR = RGBColor(0x1E, 0xB7, 0x42)  # Green for complete
PENDING_COLOR = RGBColor(0xF6, 0xBB, 0x00)  # Yellow for pending


def add_title_slide(prs, title, subtitle):
    """Add a title slide with gradient-like styling."""
    slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)

    # Add background shape
    bg_shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height
    )
    bg_shape.fill.solid()
    bg_shape.fill.fore_color.rgb = PRIMARY_COLOR
    bg_shape.line.fill.background()

    # Add title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(12.333), Inches(1.5))
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0xff, 0xff, 0xff)
    p.alignment = PP_ALIGN.CENTER

    # Add subtitle
    subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(12.333), Inches(1))
    tf = subtitle_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = subtitle
    p.font.size = Pt(24)
    p.font.color.rgb = RGBColor(0xbb, 0xe3, 0xf9)  # Light cyan complement to blue
    p.alignment = PP_ALIGN.CENTER

    return slide


def add_section_slide(prs, title, section_number=None):
    """Add a section divider slide."""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # Add colored bar on left
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(0.3), prs.slide_height)
    bar.fill.solid()
    bar.fill.fore_color.rgb = SECONDARY_COLOR
    bar.line.fill.background()

    # Section number circle
    if section_number:
        circle = slide.shapes.add_shape(
            MSO_SHAPE.OVAL, Inches(1), Inches(2.8), Inches(1), Inches(1)
        )
        circle.fill.solid()
        circle.fill.fore_color.rgb = SECONDARY_COLOR
        circle.line.fill.background()

        num_box = slide.shapes.add_textbox(Inches(1), Inches(2.95), Inches(1), Inches(0.7))
        tf = num_box.text_frame
        p = tf.paragraphs[0]
        p.text = str(section_number)
        p.font.size = Pt(36)
        p.font.bold = True
        p.font.color.rgb = RGBColor(0xff, 0xff, 0xff)
        p.alignment = PP_ALIGN.CENTER

    # Title
    title_box = slide.shapes.add_textbox(Inches(2.5), Inches(2.7), Inches(10), Inches(1.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = DARK_TEXT

    return slide


def add_content_slide(prs, title, content_items, subtitle=None):
    """Add a content slide with bullet points."""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # Header bar
    header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.2))
    header.fill.solid()
    header.fill.fore_color.rgb = PRIMARY_COLOR
    header.line.fill.background()

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0xff, 0xff, 0xff)

    # Subtitle if provided
    start_y = Inches(1.5)
    if subtitle:
        sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.4), Inches(12), Inches(0.5))
        tf = sub_box.text_frame
        p = tf.paragraphs[0]
        p.text = subtitle
        p.font.size = Pt(18)
        p.font.italic = True
        p.font.color.rgb = LIGHT_TEXT
        start_y = Inches(2)

    # Content
    content_box = slide.shapes.add_textbox(Inches(0.5), start_y, Inches(12.333), Inches(5))
    tf = content_box.text_frame
    tf.word_wrap = True

    for i, item in enumerate(content_items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()

        if isinstance(item, tuple):
            # (text, level)
            text, level = item
            p.text = text
            p.level = level
            p.font.size = Pt(18 - level * 2)
        else:
            p.text = f"• {item}"
            p.font.size = Pt(20)

        p.font.color.rgb = DARK_TEXT
        p.space_after = Pt(8)

    return slide


def add_table_slide(prs, title, headers, rows, subtitle=None):
    """Add a slide with a table."""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # Header bar
    header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.2))
    header.fill.solid()
    header.fill.fore_color.rgb = PRIMARY_COLOR
    header.line.fill.background()

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0xff, 0xff, 0xff)

    # Subtitle
    table_top = Inches(1.5)
    if subtitle:
        sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.4), Inches(12), Inches(0.4))
        tf = sub_box.text_frame
        p = tf.paragraphs[0]
        p.text = subtitle
        p.font.size = Pt(16)
        p.font.italic = True
        p.font.color.rgb = LIGHT_TEXT
        table_top = Inches(1.9)

    # Create table
    num_cols = len(headers)
    num_rows = len(rows) + 1

    table_width = Inches(12.333)
    table_height = Inches(0.5) * num_rows

    table = slide.shapes.add_table(
        num_rows, num_cols, Inches(0.5), table_top, table_width, table_height
    ).table

    # Set column widths
    col_width = table_width / num_cols
    for i in range(num_cols):
        table.columns[i].width = int(col_width)

    # Header row
    for i, header_text in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = header_text
        cell.fill.solid()
        cell.fill.fore_color.rgb = SECONDARY_COLOR
        para = cell.text_frame.paragraphs[0]
        para.font.bold = True
        para.font.size = Pt(14)
        para.font.color.rgb = RGBColor(0xff, 0xff, 0xff)
        para.alignment = PP_ALIGN.CENTER
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE

    # Data rows
    for row_idx, row_data in enumerate(rows):
        for col_idx, cell_text in enumerate(row_data):
            cell = table.cell(row_idx + 1, col_idx)
            cell.text = str(cell_text)
            para = cell.text_frame.paragraphs[0]
            para.font.size = Pt(12)
            para.font.color.rgb = DARK_TEXT
            para.alignment = PP_ALIGN.LEFT if col_idx == 0 else PP_ALIGN.CENTER
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE

            # Alternate row colors
            if row_idx % 2 == 1:
                cell.fill.solid()
                cell.fill.fore_color.rgb = LIGHT_BG

    return slide


def add_two_column_slide(prs, title, left_title, left_items, right_title, right_items):
    """Add a two-column content slide."""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # Header bar
    header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.2))
    header.fill.solid()
    header.fill.fore_color.rgb = PRIMARY_COLOR
    header.line.fill.background()

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0xff, 0xff, 0xff)

    # Left column title
    left_title_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(6), Inches(0.5))
    tf = left_title_box.text_frame
    p = tf.paragraphs[0]
    p.text = left_title
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = SECONDARY_COLOR

    # Left column content
    left_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.1), Inches(5.8), Inches(4.5))
    tf = left_box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(left_items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f"• {item}"
        p.font.size = Pt(16)
        p.font.color.rgb = DARK_TEXT
        p.space_after = Pt(6)

    # Right column title
    right_title_box = slide.shapes.add_textbox(Inches(6.8), Inches(1.5), Inches(6), Inches(0.5))
    tf = right_title_box.text_frame
    p = tf.paragraphs[0]
    p.text = right_title
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = ACCENT_COLOR

    # Right column content
    right_box = slide.shapes.add_textbox(Inches(6.8), Inches(2.1), Inches(5.8), Inches(4.5))
    tf = right_box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(right_items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f"• {item}"
        p.font.size = Pt(16)
        p.font.color.rgb = DARK_TEXT
        p.space_after = Pt(6)

    return slide


# ============================================================================
# CREATE SLIDES
# ============================================================================

# Slide 1: Title
add_title_slide(
    prs,
    "Autism Pathway Framework",
    "Technical Architecture & Machine Learning Algorithms"
)

# Slide 2: Agenda
add_content_slide(prs, "Presentation Overview", [
    "Project Vision & Objectives",
    "Core Technology Stack",
    "Data Infrastructure (Phase 1)",
    "Implemented ML Algorithms",
    "Knowledge Graph Architecture",
    "Planned ML Algorithms (Phases 2-5)",
    "Implementation Roadmap & Status"
])

# Slide 3: Section - Project Overview
add_section_slide(prs, "Project Vision & Objectives", 1)

# Slide 4: Project Overview
add_content_slide(
    prs,
    "Autism Pathway Framework",
    [
        "Multi-layer computational framework for autism genetics research",
        "Integrates genomic variants with biological knowledge graphs",
        "Enables patient stratification through pathway-based analysis",
        "Combines neural networks with symbolic reasoning",
        "Generates therapeutic hypotheses through causal inference",
        "Designed for reproducibility and clinical translation"
    ],
    subtitle="A comprehensive bioinformatics pipeline for understanding autism spectrum disorder genetics"
)

# Slide 5: Section - Tech Stack
add_section_slide(prs, "Core Technology Stack", 2)

# Slide 6: Tech Stack Overview
add_table_slide(
    prs,
    "Technology Stack Overview",
    ["Category", "Technologies", "Purpose"],
    [
        ["Language", "Python 3.10+", "Core development"],
        ["Data Science", "numpy, pandas, scipy, scikit-learn", "Numerical computing & ML"],
        ["Bioinformatics", "pysam, pyvcf3, biopython, obonet", "Genomic data processing"],
        ["Graph/Network", "NetworkX, igraph", "Knowledge graph construction"],
        ["Deep Learning", "PyTorch, DGL, torch-geometric", "Neural network models"],
        ["NLP/Embeddings", "transformers, Geneformer, ESM-2", "Pretrained biological models"],
        ["Visualization", "matplotlib, seaborn, plotly", "Data visualization"],
        ["Cloud (optional)", "GCP: Storage, BigQuery, AI Platform", "Scalable deployment"],
    ]
)

# Slide 7: Data Stack Details
add_two_column_slide(
    prs,
    "Data Infrastructure",
    "Data Formats",
    [
        "VCF: Variant Call Format files",
        "Parquet: Columnar storage (pyarrow)",
        "HDF5: Large matrix storage (h5py)",
        "AnnData: Single-cell data (.h5ad)",
        "OBO: Gene Ontology format",
        "GMT: Gene set collections",
        "YAML: Configuration files"
    ],
    "External Databases",
    [
        "gnomAD: Population frequencies",
        "ClinVar: Clinical variant annotations",
        "Gene Ontology: Functional annotations",
        "KEGG/Reactome: Pathway databases",
        "STRING: Protein interactions",
        "BrainSpan: Developmental expression",
        "SFARI Gene: Autism gene database"
    ]
)

# Slide 8: Development Stack
add_two_column_slide(
    prs,
    "Development & DevOps Stack",
    "Testing & Quality",
    [
        "pytest: Unit and integration testing",
        "pytest-cov: Code coverage",
        "pytest-mock: Mocking framework",
        "mypy: Static type checking",
        "black: Code formatting",
        "ruff: Fast Python linting"
    ],
    "Development Tools",
    [
        "click: CLI framework",
        "rich: Terminal formatting",
        "tqdm: Progress bars",
        "python-dotenv: Environment config",
        "ipython/jupyter: Interactive analysis",
        "Git: Version control"
    ]
)

# Slide 9: Section - Implemented
add_section_slide(prs, "Implemented ML Algorithms", 3)

# Slide 10: Module 01 - Data Loaders
add_content_slide(
    prs,
    "Module 01: Data Loaders",
    [
        "VCF Loader — Parse variant call format files with genotype extraction",
        "Annotation Loader — Load gene annotations from external databases",
        "Pathway Loader — Import GO, KEGG, Reactome, custom GMT pathways",
        "Expression Loader — BrainSpan developmental expression data",
        "Single-Cell Loader — Allen Brain Atlas cell-type expression",
        "Constraint Loader — gnomAD pLI/LOEUF scores, SFARI gene database"
    ],
    subtitle="Phase 1A: Data Foundation — Status: ✅ COMPLETE"
)

# Slide 11: Module 02 - Variant Processing
add_content_slide(
    prs,
    "Module 02: Variant Processing",
    [
        "Variant Annotation — Consequence classification (LoF, missense, synonymous)",
        ("Impact levels: HIGH, MODERATE, LOW, MODIFIER", 1),
        ("Pathogenicity scoring: CADD, REVEL, PolyPhen, SIFT", 1),
        "Quality Control Filters — Threshold-based variant filtering",
        ("Quality scores, allele frequency, genotype quality", 1),
        "Gene Burden Calculator — Weighted variant-to-gene aggregation",
        ("Weighting: consequence-based, CADD, allele frequency", 1),
        ("Aggregation: weighted_sum, max, count", 1),
        ("Normalization: z-score, minmax, rank", 1)
    ],
    subtitle="Phase 1B: Variant Processing — Status: ✅ COMPLETE"
)

# Slide 12: Gene Burden Algorithm
add_content_slide(
    prs,
    "Gene Burden Algorithm Details",
    [
        "Core Formula: burden_score = Σ(variant_weight × impact_weight × af_weight)",
        "Consequence Weights:",
        ("Loss-of-function (LoF): 1.0", 1),
        ("Damaging missense: 0.5", 1),
        ("Other missense: 0.3", 1),
        ("Synonymous: 0.1", 1),
        "Allele Frequency Weight: (1 - AF)^β where β controls rare variant emphasis",
        "Specialized Methods:",
        ("compute_lof_burden() — LoF variants only", 1),
        ("compute_missense_burden() — Damaging missense only", 1),
        ("combine_burdens() — Multi-type combination", 1)
    ],
    subtitle="Variant → Gene aggregation with biological priors"
)

# Slide 13: Module 03 - Knowledge Graph
add_content_slide(
    prs,
    "Module 03: Knowledge Graph",
    [
        "Graph Schema Definition — Typed nodes and edges",
        ("Node types: GENE, PATHWAY, GO_TERM, CELL_TYPE, DRUG", 1),
        ("Edge types: GENE_IN_PATHWAY, GENE_INTERACTS, GO_IS_A, etc.", 1),
        "Knowledge Graph Builder — Heterogeneous graph construction",
        ("Integrates pathways, PPI networks, GO hierarchy", 1),
        ("Backend: NetworkX MultiDiGraph", 1),
        "Graph Exporters — Multiple output formats",
        ("Neo4j Cypher export", 1),
        ("DGL heterogeneous graph", 1),
        ("PyTorch Geometric HeteroData", 1)
    ],
    subtitle="Phase 2A: Knowledge Representation — Status: ✅ COMPLETE"
)

# Slide 14: Section - Planned
add_section_slide(prs, "Planned ML Algorithms", 4)

# Slide 15: Phase 2 - Embeddings
add_content_slide(
    prs,
    "Phase 2: Graph & Pretrained Embeddings",
    [
        "Module 04: Graph Embeddings",
        ("TransE — Translational embeddings: h + r ≈ t", 1),
        ("RotatE — Rotational embeddings for complex relations", 1),
        "Module 05: Pretrained Biological Embeddings",
        ("Geneformer — Gene representations from single-cell data", 1),
        ("ESM-2 — Protein sequence embeddings for variant effects", 1),
        ("PubMedBERT/BioGPT — Literature-derived gene functions", 1),
        "Embedding Fusion Strategies",
        ("Concatenation, weighted sum, cross-attention", 1),
        ("Multi-task fine-tuning for autism-specific tasks", 1)
    ],
    subtitle="Modules 04-05 — Status: 🔲 NOT STARTED"
)

# Slide 16: Phase 3A - GNN
add_content_slide(
    prs,
    "Phase 3A: Ontology-Aware Graph Neural Network",
    [
        "Heterogeneous GNN Architecture",
        ("Edge-type specific transformations", 1),
        ("Multi-layer message passing", 1),
        "Biological Attention Mechanism",
        ("Attention weighted by developmental stage importance", 1),
        ("Cell-type specificity priors", 1),
        ("Gene constraint score integration", 1),
        "Hierarchical Aggregator",
        ("Follows GO hierarchy (is-a, part-of)", 1),
        ("Multi-level pathway aggregation", 1)
    ],
    subtitle="Module 06 — Status: 🔲 NOT STARTED"
)

# Slide 17: Phase 3B - Pathway Scoring
add_content_slide(
    prs,
    "Phase 3B: Pathway Scoring & Network Refinement",
    [
        "Gene → Pathway Aggregation",
        ("Raw score: Σ(gene_burden × gene_weight)", 1),
        ("Size normalization: score / √(pathway_size)", 1),
        ("GSEA-like enrichment scoring", 1),
        "Gene Weighting Schemes",
        ("Constraint-based (pLI, LOEUF)", 1),
        ("Expression-based (brain specificity)", 1),
        ("Network centrality", 1),
        "Random Walk with Restart (RWR)",
        ("Update: p(t+1) = (1-α)×W×p(t) + α×p(0)", 1),
        ("Hub correction with degree penalty", 1)
    ],
    subtitle="Module 07 — Status: 🔲 NOT STARTED"
)

# Slide 18: Phase 3C - Clustering
add_content_slide(
    prs,
    "Phase 3C: Subtype Clustering",
    [
        "Feature Preparation",
        ("Pathway profile matrix (samples × pathways)", 1),
        ("PCA dimensionality reduction (90% variance)", 1),
        "Clustering Algorithms",
        ("Gaussian Mixture Model (GMM) — Soft probabilistic clusters", 1),
        ("Spectral Clustering — Non-convex cluster shapes", 1),
        ("Hierarchical Clustering — Ward linkage dendrograms", 1),
        "Stability Assessment",
        ("Bootstrap resampling (100+ iterations)", 1),
        ("Co-clustering matrix analysis", 1),
        ("Stability threshold: 0.7", 1)
    ],
    subtitle="Module 08 — Status: 🔲 NOT STARTED"
)

# Slide 19: Phase 4 - Symbolic
add_content_slide(
    prs,
    "Phase 4: Symbolic Reasoning & Integration",
    [
        "Module 09: Curated Biological Rules",
        ("R1: LoF + constrained gene + cortex expression → high-confidence", 1),
        ("R2: Multiple pathway hits (≥2 genes) → convergence", 1),
        ("R3: CHD8 cascade → chromatin regulation disruption", 1),
        ("R4: Synaptic gene + excitatory neuron → synaptic subtype", 1),
        ("R5: Paralog intact + expressed → compensation", 1),
        ("R6: Drug targets pathway → therapeutic hypothesis", 1),
        "Module 10: Neuro-Symbolic Integration",
        ("Combine GNN predictions + rule-based inference", 1),
        ("Explainable reasoning chains", 1)
    ],
    subtitle="Modules 09-10 — Status: 🔲 NOT STARTED"
)

# Slide 20: Phase 5 - Causal
add_content_slide(
    prs,
    "Phase 5: Causal Inference",
    [
        "Structural Causal Model (SCM)",
        ("Causal chain: Variants → Gene → Pathway → Circuit → Phenotype", 1),
        ("Confounder adjustment: ancestry, batch effects", 1),
        "Do-Calculus (Pearl's Intervention Formalism)",
        ("Intervention queries: P(outcome | do(treatment))", 1),
        ("Average Treatment Effect (ATE)", 1),
        ("Conditional ATE (CATE) for subgroups", 1),
        "Counterfactual Reasoning",
        ("Three-step: Abduction → Action → Prediction", 1),
        ("Probability of Necessity/Sufficiency", 1),
        "Mediation Analysis",
        ("Natural Direct/Indirect Effects", 1)
    ],
    subtitle="Module 12 — Status: 🔲 NOT STARTED"
)

# Slide 21: Section - Roadmap
add_section_slide(prs, "Implementation Roadmap", 5)

# Slide 22: Progress Table
add_table_slide(
    prs,
    "Implementation Status",
    ["Phase", "Modules", "Description", "Status", "Progress"],
    [
        ["Phase 1", "01-02", "Data Foundation", "✅ Complete", "100%"],
        ["Phase 2A", "03", "Knowledge Graph", "✅ Complete", "100%"],
        ["Phase 2B-C", "04-05", "Embeddings", "🔲 Not Started", "0%"],
        ["Phase 3", "06-08", "Neural Models", "🔲 Not Started", "0%"],
        ["Phase 4", "09-11", "Symbolic Reasoning", "🔲 Not Started", "0%"],
        ["Phase 5", "12", "Causal Inference", "🔲 Not Started", "0%"],
    ],
    subtitle="Overall Progress: ~25% (3 of 12 modules complete)"
)

# Slide 23: Architecture Diagram (Text representation)
add_content_slide(
    prs,
    "System Architecture Overview",
    [
        "Layer 1: Data Ingestion",
        ("VCF → Variants → Annotated Variants → Gene Burden Matrix", 1),
        "Layer 2: Knowledge Representation",
        ("Knowledge Graph (genes, pathways, GO, PPI)", 1),
        ("Graph embeddings + Pretrained embeddings", 1),
        "Layer 3: Neural Processing",
        ("Ontology-aware GNN → Pathway scores", 1),
        ("Subtype clustering → Patient stratification", 1),
        "Layer 4: Reasoning & Interpretation",
        ("Symbolic rules + Neural predictions → Explanations", 1),
        ("Causal inference → Therapeutic hypotheses", 1)
    ],
    subtitle="End-to-end pipeline from variants to actionable insights"
)

# Slide 24: Key Differentiators
add_two_column_slide(
    prs,
    "Framework Differentiators",
    "Technical Innovation",
    [
        "Heterogeneous knowledge graph",
        "Ontology-aware GNN architecture",
        "Multi-modal embedding fusion",
        "Neuro-symbolic integration",
        "Causal inference framework",
        "Reproducible analysis pipeline"
    ],
    "Biological Relevance",
    [
        "Brain-specific expression context",
        "Developmental stage awareness",
        "Cell-type specificity",
        "Constraint-aware scoring",
        "Pathway convergence detection",
        "Therapeutic target prioritization"
    ]
)

# Slide 25: Closing
add_title_slide(
    prs,
    "Thank You",
    "Autism Pathway Framework — Bridging Genomics and Clinical Insights"
)

# Save presentation
output_path = "Autism_Pathway_Framework_Tech_Stack.pptx"
prs.save(output_path)
print(f"Presentation saved to: {output_path}")
print(f"Total slides: {len(prs.slides)}")
