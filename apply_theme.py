#!/usr/bin/env python3
"""
Apply Autism Awareness color scheme to AI_Genetics_Project_Intro.pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.dml import MSO_THEME_COLOR
import os

# Color scheme - Autism Awareness Theme
PRIMARY_COLOR = RGBColor(0x0C, 0x96, 0xE4)  # #0C96E4 - Blue Cola
SECONDARY_COLOR = RGBColor(0x1E, 0xB7, 0x42)  # #1EB742 - American Green
ACCENT_COLOR = RGBColor(0xFF, 0x33, 0x34)  # #FF3334 - Deep Carmine Pink
HIGHLIGHT_COLOR = RGBColor(0xF6, 0xBB, 0x00)  # #F6BB00 - American Yellow
DARK_TEXT = RGBColor(0x1f, 0x29, 0x37)  # Dark slate
LIGHT_TEXT = RGBColor(0x4b, 0x55, 0x63)  # Medium gray
WHITE = RGBColor(0xff, 0xff, 0xff)
LIGHT_BG = RGBColor(0xf0, 0xf9, 0xff)  # Light blue tint

# Input file
input_file = "/Users/rohitchauhan/Downloads/AI-Genetic-Research/autism-pathway-framework/AI_Genetics_Project_Intro.pptx"
output_file = "/Users/rohitchauhan/Downloads/AI-Genetic-Research/autism-pathway-framework/AI_Genetics_Project_Intro_Themed.pptx"

# Load presentation
prs = Presentation(input_file)

print(f"Loaded presentation with {len(prs.slides)} slides")

def apply_color_to_shape(shape, is_title=False, is_background=False):
    """Apply theme colors to a shape."""
    try:
        # Handle shapes with fill
        if hasattr(shape, 'fill'):
            fill = shape.fill
            if fill.type is not None:
                # Check if it's a solid fill
                if hasattr(fill, 'fore_color') and fill.fore_color is not None:
                    try:
                        current_color = fill.fore_color.rgb
                        # Apply primary color to dark/blue backgrounds
                        if is_background or is_title:
                            fill.solid()
                            fill.fore_color.rgb = PRIMARY_COLOR
                    except:
                        pass
    except Exception as e:
        pass

def apply_color_to_text(paragraph, is_title=False, is_subtitle=False, on_dark_bg=False):
    """Apply theme colors to text."""
    try:
        for run in paragraph.runs:
            if run.font:
                if on_dark_bg:
                    run.font.color.rgb = WHITE
                elif is_title:
                    run.font.color.rgb = PRIMARY_COLOR
                elif is_subtitle:
                    run.font.color.rgb = LIGHT_TEXT
                else:
                    run.font.color.rgb = DARK_TEXT
    except Exception as e:
        pass

def is_dark_shape(shape):
    """Check if a shape has a dark fill color."""
    try:
        if hasattr(shape, 'fill') and shape.fill.type is not None:
            if hasattr(shape.fill, 'fore_color') and shape.fill.fore_color is not None:
                rgb = shape.fill.fore_color.rgb
                if rgb:
                    # Calculate luminance
                    r, g, b = rgb[0], rgb[1], rgb[2]
                    luminance = (0.299 * r + 0.587 * g + 0.114 * b) / 255
                    return luminance < 0.5
    except:
        pass
    return False

def process_shape(shape, slide_idx):
    """Process a single shape and apply colors."""
    shape_type = type(shape).__name__

    # Handle grouped shapes
    if hasattr(shape, 'shapes'):
        for sub_shape in shape.shapes:
            process_shape(sub_shape, slide_idx)
        return

    # Check if shape has dark background
    has_dark_bg = is_dark_shape(shape)

    # Apply fill colors to shapes
    if hasattr(shape, 'fill'):
        try:
            fill = shape.fill
            if fill.type is not None and hasattr(fill, 'fore_color'):
                if fill.fore_color is not None:
                    try:
                        rgb = fill.fore_color.rgb
                        if rgb:
                            r, g, b = rgb[0], rgb[1], rgb[2]
                            luminance = (0.299 * r + 0.587 * g + 0.114 * b) / 255

                            # Dark colors -> Primary Blue
                            if luminance < 0.3:
                                fill.solid()
                                fill.fore_color.rgb = PRIMARY_COLOR
                            # Medium dark colors -> Secondary Green
                            elif luminance < 0.5:
                                fill.solid()
                                fill.fore_color.rgb = SECONDARY_COLOR
                            # Accent colors (saturated) -> keep or use accent
                            elif r > 200 and g < 100:  # Reddish
                                fill.solid()
                                fill.fore_color.rgb = ACCENT_COLOR
                            elif g > 200 and r < 100:  # Greenish
                                fill.solid()
                                fill.fore_color.rgb = SECONDARY_COLOR
                            elif r > 200 and g > 150 and b < 100:  # Yellowish
                                fill.solid()
                                fill.fore_color.rgb = HIGHLIGHT_COLOR
                    except:
                        pass
        except:
            pass

    # Apply colors to text
    if hasattr(shape, 'text_frame'):
        try:
            tf = shape.text_frame
            for para in tf.paragraphs:
                for run in para.runs:
                    if run.font:
                        try:
                            # Check if text is on a dark background
                            if has_dark_bg:
                                run.font.color.rgb = WHITE
                            else:
                                # Check current color
                                current = run.font.color.rgb
                                if current:
                                    r, g, b = current[0], current[1], current[2]
                                    luminance = (0.299 * r + 0.587 * g + 0.114 * b) / 255

                                    # Very light text (on dark bg) -> White
                                    if luminance > 0.9:
                                        pass  # Keep white
                                    # Dark text -> Dark text color
                                    elif luminance < 0.3:
                                        run.font.color.rgb = DARK_TEXT
                                    # Colored text -> Apply theme colors
                                    else:
                                        # Blue-ish -> Primary
                                        if b > r and b > g:
                                            run.font.color.rgb = PRIMARY_COLOR
                                        # Green-ish -> Secondary
                                        elif g > r and g > b:
                                            run.font.color.rgb = SECONDARY_COLOR
                                        # Red-ish -> Accent
                                        elif r > g and r > b:
                                            run.font.color.rgb = ACCENT_COLOR
                        except:
                            pass
        except:
            pass

    # Handle tables
    if hasattr(shape, 'table'):
        try:
            table = shape.table
            for row_idx, row in enumerate(table.rows):
                for col_idx, cell in enumerate(row.cells):
                    # Header row (first row)
                    if row_idx == 0:
                        cell.fill.solid()
                        cell.fill.fore_color.rgb = SECONDARY_COLOR
                        for para in cell.text_frame.paragraphs:
                            for run in para.runs:
                                if run.font:
                                    run.font.color.rgb = WHITE
                                    run.font.bold = True
                    else:
                        # Alternating row colors
                        if row_idx % 2 == 0:
                            cell.fill.solid()
                            cell.fill.fore_color.rgb = LIGHT_BG
                        for para in cell.text_frame.paragraphs:
                            for run in para.runs:
                                if run.font:
                                    run.font.color.rgb = DARK_TEXT
        except Exception as e:
            pass

# Process each slide
for slide_idx, slide in enumerate(prs.slides):
    print(f"Processing slide {slide_idx + 1}...")

    # Process all shapes on the slide
    for shape in slide.shapes:
        process_shape(shape, slide_idx)

    # Handle slide background if accessible
    try:
        if slide.background and slide.background.fill:
            bg_fill = slide.background.fill
            if bg_fill.type is not None:
                # Don't change white/light backgrounds
                pass
    except:
        pass

# Save the themed presentation
prs.save(output_file)
print(f"\nThemed presentation saved to: {output_file}")
print("Color scheme applied:")
print(f"  Primary (Blue Cola):     #{PRIMARY_COLOR}")
print(f"  Secondary (Am. Green):   #{SECONDARY_COLOR}")
print(f"  Accent (Carmine Pink):   #{ACCENT_COLOR}")
print(f"  Highlight (Am. Yellow):  #{HIGHLIGHT_COLOR}")
