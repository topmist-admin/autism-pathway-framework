#!/usr/bin/env python3
"""
Create a merged, concise PowerPoint presentation combining:
- AI_Genetics_Project_Intro_Themed.pptx (flow/narrative)
- Autism_Pathway_Framework_Tech_Stack.pptx (technical details)

Uses Autism Awareness color scheme.
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# Create presentation with widescreen dimensions
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# ============================================================================
# COLOR SCHEME - Autism Awareness Theme
# ============================================================================
PRIMARY_COLOR = RGBColor(0x0C, 0x96, 0xE4)  # #0C96E4 - Blue Cola
SECONDARY_COLOR = RGBColor(0x1E, 0xB7, 0x42)  # #1EB742 - American Green
ACCENT_COLOR = RGBColor(0xFF, 0x33, 0x34)  # #FF3334 - Deep Carmine Pink
HIGHLIGHT_COLOR = RGBColor(0xF6, 0xBB, 0x00)  # #F6BB00 - American Yellow
DARK_TEXT = RGBColor(0x1f, 0x29, 0x37)
LIGHT_TEXT = RGBColor(0x4b, 0x55, 0x63)
WHITE = RGBColor(0xff, 0xff, 0xff)
LIGHT_BG = RGBColor(0xf0, 0xf9, 0xff)
LIGHT_CYAN = RGBColor(0xbb, 0xe3, 0xf9)

# ============================================================================
# SLIDE HELPER FUNCTIONS
# ============================================================================

def add_title_slide(title, subtitle):
    """Add a title slide with colored background."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Background
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = PRIMARY_COLOR
    bg.line.fill.background()

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(12.333), Inches(1.5))
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    # Subtitle
    sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(12.333), Inches(1))
    tf = sub_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = subtitle
    p.font.size = Pt(24)
    p.font.color.rgb = LIGHT_CYAN
    p.alignment = PP_ALIGN.CENTER

    return slide


def add_section_slide(title, section_num=None):
    """Add a section divider slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Colored bar on left
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(0.3), prs.slide_height)
    bar.fill.solid()
    bar.fill.fore_color.rgb = SECONDARY_COLOR
    bar.line.fill.background()

    # Section number circle
    if section_num:
        circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(1), Inches(2.8), Inches(1), Inches(1))
        circle.fill.solid()
        circle.fill.fore_color.rgb = SECONDARY_COLOR
        circle.line.fill.background()

        num_box = slide.shapes.add_textbox(Inches(1), Inches(2.95), Inches(1), Inches(0.7))
        tf = num_box.text_frame
        p = tf.paragraphs[0]
        p.text = str(section_num)
        p.font.size = Pt(36)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER

    # Title
    title_box = slide.shapes.add_textbox(Inches(2.5), Inches(2.7), Inches(10), Inches(1.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = DARK_TEXT

    return slide


def add_content_slide(title, items, subtitle=None):
    """Add a content slide with bullet points."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Header bar
    header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.2))
    header.fill.solid()
    header.fill.fore_color.rgb = PRIMARY_COLOR
    header.line.fill.background()

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = WHITE

    # Subtitle
    start_y = Inches(1.5)
    if subtitle:
        sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.4), Inches(12), Inches(0.5))
        tf = sub_box.text_frame
        p = tf.paragraphs[0]
        p.text = subtitle
        p.font.size = Pt(18)
        p.font.italic = True
        p.font.color.rgb = LIGHT_TEXT
        start_y = Inches(2)

    # Content
    content_box = slide.shapes.add_textbox(Inches(0.5), start_y, Inches(12.333), Inches(5))
    tf = content_box.text_frame
    tf.word_wrap = True

    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()

        if isinstance(item, tuple):
            text, level = item
            p.text = text
            p.level = level
            p.font.size = Pt(18 - level * 2)
        else:
            p.text = f"• {item}"
            p.font.size = Pt(20)

        p.font.color.rgb = DARK_TEXT
        p.space_after = Pt(8)

    return slide


def add_two_column_slide(title, left_title, left_items, right_title, right_items):
    """Add a two-column content slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Header
    header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.2))
    header.fill.solid()
    header.fill.fore_color.rgb = PRIMARY_COLOR
    header.line.fill.background()

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = WHITE

    # Left column title
    left_title_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(6), Inches(0.5))
    tf = left_title_box.text_frame
    p = tf.paragraphs[0]
    p.text = left_title
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = SECONDARY_COLOR

    # Left column content
    left_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.1), Inches(5.8), Inches(4.5))
    tf = left_box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(left_items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f"• {item}"
        p.font.size = Pt(16)
        p.font.color.rgb = DARK_TEXT
        p.space_after = Pt(6)

    # Right column title
    right_title_box = slide.shapes.add_textbox(Inches(6.8), Inches(1.5), Inches(6), Inches(0.5))
    tf = right_title_box.text_frame
    p = tf.paragraphs[0]
    p.text = right_title
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = ACCENT_COLOR

    # Right column content
    right_box = slide.shapes.add_textbox(Inches(6.8), Inches(2.1), Inches(5.8), Inches(4.5))
    tf = right_box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(right_items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f"• {item}"
        p.font.size = Pt(16)
        p.font.color.rgb = DARK_TEXT
        p.space_after = Pt(6)

    return slide


def add_table_slide(title, headers, rows, subtitle=None):
    """Add a slide with a table."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Header
    header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.2))
    header.fill.solid()
    header.fill.fore_color.rgb = PRIMARY_COLOR
    header.line.fill.background()

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = WHITE

    # Subtitle
    table_top = Inches(1.5)
    if subtitle:
        sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.4), Inches(12), Inches(0.4))
        tf = sub_box.text_frame
        p = tf.paragraphs[0]
        p.text = subtitle
        p.font.size = Pt(16)
        p.font.italic = True
        p.font.color.rgb = LIGHT_TEXT
        table_top = Inches(1.9)

    # Table
    num_cols = len(headers)
    num_rows = len(rows) + 1
    table_width = Inches(12.333)
    table_height = Inches(0.45) * num_rows

    table = slide.shapes.add_table(num_rows, num_cols, Inches(0.5), table_top, table_width, table_height).table

    col_width = table_width / num_cols
    for i in range(num_cols):
        table.columns[i].width = int(col_width)

    # Header row
    for i, header_text in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = header_text
        cell.fill.solid()
        cell.fill.fore_color.rgb = SECONDARY_COLOR
        para = cell.text_frame.paragraphs[0]
        para.font.bold = True
        para.font.size = Pt(14)
        para.font.color.rgb = WHITE
        para.alignment = PP_ALIGN.CENTER
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE

    # Data rows
    for row_idx, row_data in enumerate(rows):
        for col_idx, cell_text in enumerate(row_data):
            cell = table.cell(row_idx + 1, col_idx)
            cell.text = str(cell_text)
            para = cell.text_frame.paragraphs[0]
            para.font.size = Pt(12)
            para.font.color.rgb = DARK_TEXT
            para.alignment = PP_ALIGN.LEFT if col_idx == 0 else PP_ALIGN.CENTER
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE

            if row_idx % 2 == 1:
                cell.fill.solid()
                cell.fill.fore_color.rgb = LIGHT_BG

    return slide


# ============================================================================
# CREATE MERGED PRESENTATION
# ============================================================================

# ----- SLIDE 1: Title -----
add_title_slide(
    "AI + Genetics: Autism Pathway Framework",
    "A Systems Engineering Approach to Understanding Complex Biological Systems"
)

# ----- SLIDE 2: The Problem -----
add_section_slide("The Core Problem", 1)

# ----- SLIDE 3: Problem Details -----
add_content_slide(
    "Why Traditional Approaches Fail",
    [
        "Some problems are too complex for single variables",
        ("Many inputs with weak individual signals", 1),
        ("Signal-to-noise ratio too low for simple methods", 1),
        "Single-gene thinking ≈ single log line debugging",
        ("Works sometimes, fails at scale", 1),
        ("Misses system-level interactions", 1),
        "Analogy: Microservices, distributed systems, network outages",
        ("Rare bugs in different components cause failures when combined", 1)
    ],
    subtitle="Complex systems require systems-level analysis"
)

# ----- SLIDE 4: Our Approach -----
add_section_slide("Our Approach", 2)

# ----- SLIDE 5: Pipeline Overview -----
add_content_slide(
    "High-Level Pipeline",
    [
        "Translate raw signals → system-level stress patterns",
        "Pipeline: Variants → Genes → Pathways → Subtypes",
        ("Variant Processing: Quality control, annotation, burden calculation", 1),
        ("Knowledge Graph: Integrate biological networks and ontologies", 1),
        ("ML Analysis: Clustering, pattern recognition, subtype discovery", 1),
        "Key Concepts:",
        ("Pathway ≈ Functional subsystem (like auth flow, messaging pipeline)", 1),
        ("Network refinement: Not all components are equal, some failures propagate", 1)
    ],
    subtitle="Systems engineering principles applied to biology"
)

# ----- SLIDE 6: Tech Stack -----
add_section_slide("Technology Stack", 3)

# ----- SLIDE 7: Tech Stack Table -----
add_table_slide(
    "Core Technologies",
    ["Category", "Technologies", "Purpose"],
    [
        ["Language", "Python 3.10+", "Core development"],
        ["Data Science", "numpy, pandas, scipy, scikit-learn", "Numerical computing & ML"],
        ["Bioinformatics", "pysam, pyvcf3, biopython, obonet", "Genomic data processing"],
        ["Graph/Network", "NetworkX, igraph", "Knowledge graph construction"],
        ["Deep Learning", "PyTorch, DGL, torch-geometric", "Neural network models"],
        ["NLP/Embeddings", "transformers, Geneformer, ESM-2", "Pretrained biological models"],
        ["Visualization", "matplotlib, seaborn, plotly", "Data visualization"],
    ]
)

# ----- SLIDE 8: Data Infrastructure -----
add_two_column_slide(
    "Data Infrastructure",
    "Data Formats",
    [
        "VCF: Variant Call Format files",
        "Parquet: Columnar storage (pyarrow)",
        "HDF5: Large matrix storage (h5py)",
        "AnnData: Single-cell data (.h5ad)",
        "OBO: Gene Ontology format",
        "GMT: Gene set collections"
    ],
    "External Databases",
    [
        "gnomAD: Population frequencies",
        "ClinVar: Clinical annotations",
        "Gene Ontology: Functional annotations",
        "KEGG/Reactome: Pathway databases",
        "STRING: Protein interactions",
        "SFARI Gene: Autism gene database"
    ]
)

# ----- SLIDE 9: Implemented Algorithms -----
add_section_slide("Implemented ML Algorithms", 4)

# ----- SLIDE 10: Data Loaders & Variant Processing -----
add_content_slide(
    "Phase 1: Data Foundation (Complete)",
    [
        "Module 01: Data Loaders",
        ("VCF Loader — Parse variants with genotype extraction", 1),
        ("Pathway Loader — GO, KEGG, Reactome, custom GMT", 1),
        ("Constraint Loader — gnomAD pLI/LOEUF, SFARI genes", 1),
        "Module 02: Variant Processing",
        ("Annotation — Consequence classification (LoF, missense, synonymous)", 1),
        ("QC Filters — Quality scores, allele frequency, genotype quality", 1),
        ("Gene Burden — Weighted variant-to-gene aggregation", 1),
        "Gene Burden Formula:",
        ("burden_score = Σ(variant_weight × impact_weight × af_weight)", 1)
    ],
    subtitle="Status: ✅ COMPLETE"
)

# ----- SLIDE 11: Knowledge Graph -----
add_content_slide(
    "Phase 2A: Knowledge Graph (Complete)",
    [
        "Module 03: Knowledge Graph Construction",
        ("Graph Schema — Typed nodes and edges", 1),
        ("Node types: GENE, PATHWAY, GO_TERM, CELL_TYPE, DRUG", 1),
        ("Edge types: GENE_IN_PATHWAY, GENE_INTERACTS, GO_IS_A", 1),
        "Knowledge Graph Builder",
        ("Integrates pathways, PPI networks, GO hierarchy", 1),
        ("Backend: NetworkX MultiDiGraph", 1),
        "Exporters",
        ("Neo4j Cypher, DGL heterogeneous graph, PyTorch Geometric", 1)
    ],
    subtitle="Status: ✅ COMPLETE"
)

# ----- SLIDE 12: Planned Algorithms -----
add_section_slide("Planned ML Algorithms", 5)

# ----- SLIDE 13: Planned Algorithms Table -----
add_table_slide(
    "Planned ML Algorithms (Modules 04-12)",
    ["Module", "Algorithm", "Description"],
    [
        ["04: Graph Embeddings", "TransE, RotatE", "Knowledge graph embeddings"],
        ["05: Pretrained", "Geneformer, ESM-2", "Gene/protein embeddings"],
        ["06: Ontology GNN", "Heterogeneous GNN", "Biological attention mechanism"],
        ["07: Pathway Scoring", "GSEA-like, RWR", "Gene→pathway aggregation"],
        ["08: Clustering", "GMM, Spectral", "Pathway-based stratification"],
        ["09-10: Symbolic", "Rule Engine", "Neuro-symbolic integration"],
        ["12: Causal", "SCM, Do-Calculus", "Causal effect estimation"],
    ],
    subtitle="Status: 🔲 NOT STARTED"
)

# ----- SLIDE 14: Implementation Status -----
add_table_slide(
    "Implementation Roadmap",
    ["Phase", "Modules", "Description", "Status"],
    [
        ["Phase 1", "01-02", "Data Foundation", "✅ Complete"],
        ["Phase 2A", "03", "Knowledge Graph", "✅ Complete"],
        ["Phase 2B-C", "04-05", "Embeddings", "🔲 Not Started"],
        ["Phase 3", "06-08", "Neural Models", "🔲 Not Started"],
        ["Phase 4", "09-11", "Symbolic Reasoning", "🔲 Not Started"],
        ["Phase 5", "12", "Causal Inference", "🔲 Not Started"],
    ],
    subtitle="Overall Progress: ~25% (3 of 12 modules complete)"
)

# ----- SLIDE 15: Validation -----
add_section_slide("Validation & Outputs", 6)

# ----- SLIDE 16: Stability & Validation -----
add_two_column_slide(
    "Stability & Validation",
    "Key Questions",
    [
        "Does the result persist if data changes?",
        "Does it replicate in other cohorts?",
        "Bootstrap resampling (100+ iterations)",
        "Co-clustering matrix analysis",
        "Stability threshold: 0.7"
    ],
    "What This Is NOT",
    [
        "Not a medical tool",
        "Not diagnostic or predictive",
        "Not for individual patients",
        "Strictly research-only",
        "No clinical recommendations"
    ]
)

# ----- SLIDE 17: Project Outputs -----
add_two_column_slide(
    "Project Outputs & Resources",
    "Deliverables",
    [
        "Research paper (framework)",
        "GitHub repository",
        "Substack essays",
        "All research-only outputs"
    ],
    "GitHub Repository",
    [
        "github.com/topmist-admin/autism-pathway-framework",
        "Mostly Python code",
        "Modular architecture",
        "Comprehensive documentation",
        "Test fixtures included"
    ]
)

# ----- SLIDE 18: Key Takeaway -----
add_content_slide(
    "Key Takeaways",
    [
        "This is a systems engineering problem applied to biology",
        ("Same principles as debugging distributed systems", 1),
        ("Different domain, universal approach", 1),
        "For: Architects, ML engineers, researchers, systems thinkers",
        "Technical Innovation:",
        ("Heterogeneous knowledge graph", 1),
        ("Ontology-aware GNN architecture", 1),
        ("Neuro-symbolic integration", 1),
        "Biological Relevance:",
        ("Brain-specific expression context", 1),
        ("Developmental stage awareness", 1),
        ("Therapeutic target prioritization", 1)
    ]
)

# ----- SLIDE 19: Thank You -----
add_title_slide(
    "Thank You",
    "AI + Genetics: Bridging Systems Engineering and Biological Research"
)

# ============================================================================
# SAVE PRESENTATION
# ============================================================================
output_path = "AI_Genetics_Framework_Merged.pptx"
prs.save(output_path)
print(f"Merged presentation saved to: {output_path}")
print(f"Total slides: {len(prs.slides)}")
