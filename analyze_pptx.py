#!/usr/bin/env python3
"""
Analyze the structure of both PowerPoint presentations.
"""

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

def analyze_presentation(filepath):
    """Analyze a presentation and extract slide information."""
    prs = Presentation(filepath)
    print(f"\n{'='*60}")
    print(f"FILE: {filepath}")
    print(f"{'='*60}")
    print(f"Total slides: {len(prs.slides)}")
    print(f"Slide dimensions: {prs.slide_width.inches:.2f}\" x {prs.slide_height.inches:.2f}\"")
    print()

    for idx, slide in enumerate(prs.slides):
        print(f"\n--- Slide {idx + 1} ---")

        # Try to find title
        title_text = ""
        subtitle_text = ""
        body_texts = []

        for shape in slide.shapes:
            if hasattr(shape, 'text') and shape.text.strip():
                text = shape.text.strip()[:100]  # First 100 chars

                # Check if it's likely a title (usually first text, larger)
                if hasattr(shape, 'text_frame'):
                    for para in shape.text_frame.paragraphs:
                        for run in para.runs:
                            if run.font.size and run.font.size.pt >= 28:
                                if not title_text:
                                    title_text = text
                                break

                if not title_text and len(text) < 80:
                    title_text = text
                elif text != title_text:
                    body_texts.append(text[:80] + "..." if len(text) > 80 else text)

        if title_text:
            print(f"  Title: {title_text}")
        if body_texts:
            print(f"  Content snippets:")
            for bt in body_texts[:3]:  # First 3 content items
                print(f"    - {bt}")

        # Count shape types
        shape_types = {}
        for shape in slide.shapes:
            stype = str(shape.shape_type).replace('MSO_SHAPE_TYPE.', '')
            shape_types[stype] = shape_types.get(stype, 0) + 1

        print(f"  Shapes: {shape_types}")

    return prs

# Analyze both presentations
print("\n" + "="*80)
print("PRESENTATION ANALYSIS")
print("="*80)

prs1 = analyze_presentation("AI_Genetics_Project_Intro_Themed.pptx")
prs2 = analyze_presentation("Autism_Pathway_Framework_Tech_Stack.pptx")
